from .base_agent import BaseAgent
import os
import base64
import io
from groq import Groq
from pypdf import PdfReader
from PIL import Image

class MultimodalAgent(BaseAgent):
    def __init__(self):
        super().__init__("MultimodalAgent")
        self.groq_key = os.getenv("GROQ_API_KEY")
        self.groq_client = None
        self.source_diagrams = []
        if self.groq_key and self.groq_key != "gsk-placeholder":
             self.groq_client = Groq(api_key=self.groq_key)

    def process(self, data):
        """
        Data expectations: {'file_path': str, 'file_type': str}
        """
        file_path = data.get('file_path')
        file_type = data.get('file_type', 'text')

        if not file_path or not os.path.exists(file_path):
            raise ValueError(f"File not found: {file_path}")

        try:
            # Handle Images (Vision)
            if file_type.startswith('image/'):
                print(f"MultimodalAgent: Processing Image {file_path}")
                self.source_diagrams.extend(self._extract_image_as_diagram(file_path))
                if self.groq_client:
                    return self._process_image(file_path)
                else:
                    return "[Image Content - No API Key Provided]"

            # Handle PDF
            elif file_type == 'application/pdf':
                print(f"MultimodalAgent: Processing PDF {file_path}")
                self.source_diagrams.extend(self._extract_pdf_images(file_path))
                return self._extract_pdf_text(file_path)

            # Handle PowerPoint
            elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                print(f"MultimodalAgent: Processing PPT {file_path}")
                self.source_diagrams.extend(self._extract_pptx_images(file_path))
                return self._extract_pptx_text(file_path)

            # Handle Audio/Video
            elif file_type.startswith('audio/') or file_type.startswith('video/'):
                print(f"MultimodalAgent: Processing Audio/Video {file_path}")
                if file_type.startswith('video/'):
                    self.source_diagrams.extend(self._extract_video_frame(file_path))
                return self._transcribe_media(file_path)
            
            # Handle Text
            else:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return content
        except Exception as e:
            print(f"Error processing file: {e}")
            return ""

    def _extract_pdf_text(self, pdf_path):
        try:
            reader = PdfReader(pdf_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text if text.strip() else "[Empty PDF]"
        except Exception as e:
            print(f"PDF Error: {e}")
            return "[Error reading PDF]"

    def _extract_pdf_images(self, pdf_path):
        try:
            import fitz  # PyMuPDF
        except Exception:
            print("PDF Image Extraction: PyMuPDF not installed")
            return []

        diagrams = []
        try:
            doc = fitz.open(pdf_path)
            for page_index in range(len(doc)):
                page = doc[page_index]
                page_title = self._get_pdf_page_title(page)
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base = doc.extract_image(xref)
                    image_bytes = base.get("image")
                    ext = base.get("ext", "png")
                    if not image_bytes:
                        continue
                    data_url = self._encode_image_bytes(image_bytes, f"image/{ext}")
                    title_suffix = f" - {page_title}" if page_title else ""
                    diagrams.append({
                        "title": f"{os.path.basename(pdf_path)} - Page {page_index + 1}{title_suffix}",
                        "data_url": data_url
                    })
                    print(f"PDF Image Extracted: Page {page_index + 1} ({len(image_bytes)} bytes) - Total: {len(diagrams)}")
        except Exception as e:
            print(f"PDF Image Extraction Error: {e}")
        print(f"PDF Complete: {len(diagrams)} images extracted")
        return diagrams

    def _extract_pptx_text(self, ppt_path):
        try:
            from pptx import Presentation
            prs = Presentation(ppt_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"PPT Error: {e}")
            return "[Error reading PowerPoint]"

    def _extract_pptx_images(self, ppt_path):
        diagrams = []
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            prs = Presentation(ppt_path)
            for slide_index, slide in enumerate(prs.slides):
                slide_title = self._get_ppt_slide_title(slide)
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            data_url = self._encode_image_bytes(image_bytes, image.content_type)
                            title_suffix = f" - {slide_title}" if slide_title else ""
                            diagrams.append({
                                "title": f"{os.path.basename(ppt_path)} - Slide {slide_index + 1}{title_suffix}",
                                "data_url": data_url
                            })
                            print(f"PPT Image Extracted: Slide {slide_index + 1} ({len(image_bytes)} bytes) - Total: {len(diagrams)}")
                        except Exception as e:
                            print(f"PPT Image Processing Error at Slide {slide_index + 1}: {e}")
                            continue
        except Exception as e:
            print(f"PPT Image Extraction Error: {e}")
        print(f"PPT Complete: {len(diagrams)} images extracted")
        return diagrams

    def _transcribe_media(self, media_path):
        try:
            # For video, extract audio first using moviepy if needed, 
            # but Groq API might accept video/audio file directly depending on size.
            # For stability, passing the file object directly to Groq Whisper.
            
            with open(media_path, "rb") as file:
                transcription = self.groq_client.audio.transcriptions.create(
                    file=(os.path.basename(media_path), file.read()),
                    model="whisper-large-v3-turbo",
                    response_format="text"
                )
            return transcription
        except Exception as e:
            print(f"Whisper Error: {e}")
            return "[Error transcribing media]"

    def _extract_video_frame(self, video_path):
        diagrams = []
        try:
            from moviepy.editor import VideoFileClip
            import numpy as np
            
            clip = VideoFileClip(video_path)
            duration = clip.duration
            print(f"Video Processing: Duration = {duration:.1f}s")
            
            # Extract multiple frames throughout the video
            num_frames = min(5, max(2, int(duration / 10)))  # 1 frame every 10 seconds, min 2, max 5
            timestamps = [duration * i / (num_frames - 1) if num_frames > 1 else duration / 2 for i in range(num_frames)]
            print(f"Video Processing: Extracting {num_frames} frames at timestamps: {[f'{t:.1f}s' for t in timestamps]}")
            
            for frame_idx, timestamp in enumerate(timestamps):
                try:
                    print(f"Video: Extracting frame {frame_idx + 1}/{num_frames} at {timestamp:.1f}s")
                    frame = clip.get_frame(timestamp)
                    image = Image.fromarray(frame)
                    
                    frame_array = np.array(image, dtype=np.float32)
                    brightness = np.mean(frame_array)
                    print(f"Video: Frame @{timestamp:.1f}s brightness={brightness:.1f}")
                    
                    if brightness < 15:
                        print(f"Video: Frame @{timestamp:.1f}s skipped (too dark)")
                        continue
                    
                    buf = io.BytesIO()
                    image.save(buf, format="PNG")
                    buf.seek(0)
                    image_data = buf.getvalue()
                    data_url = self._encode_image_bytes(image_data, "image/png")
                    diagrams.append({
                        "title": f"{os.path.basename(video_path)} - Frame @ {timestamp:.1f}s",
                        "data_url": data_url
                    })
                    print(f"Video: Frame extracted @{timestamp:.1f}s - Total: {len(diagrams)}")
                except Exception as e:
                    print(f"Video: Frame error at {timestamp:.1f}s: {e}")
                    import traceback
                    traceback.print_exc()
                    continue
            
            clip.close()
            print(f"Video: Complete - {len(diagrams)} frames extracted")
        except Exception as e:
            print(f"Video: Extraction error - {e}")
            import traceback
            traceback.print_exc()
        return diagrams

    def _extract_image_as_diagram(self, image_path):
        try:
            with open(image_path, "rb") as image_file:
                data_url = self._encode_image_bytes(image_file.read(), "image/png")
            return [{
                "title": os.path.basename(image_path),
                "data_url": data_url
            }]
        except Exception as e:
            print(f"Image Extraction Error: {e}")
            return []

    def _encode_image_bytes(self, image_bytes, mime_type):
        encoded = base64.b64encode(image_bytes).decode('utf-8')
        return f"data:{mime_type};base64,{encoded}"

    def _get_pdf_page_title(self, page):
        try:
            text = page.get_text("text") or ""
            for line in text.splitlines():
                cleaned = line.strip()
                if cleaned:
                    return cleaned[:60]
        except Exception:
            return ""
        return ""

    def _get_ppt_slide_title(self, slide):
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    cleaned = shape.text.strip().splitlines()[0]
                    if cleaned:
                        return cleaned[:60]
        except Exception:
            return ""
        return ""

    def _process_image(self, image_path):
        try:
            with open(image_path, "rb") as image_file:
                base64_image = base64.b64encode(image_file.read()).decode('utf-8')

            chat_completion = self.groq_client.chat.completions.create(
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Describe this image in detail to extract key concepts and relationships for a knowledge graph."},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{base64_image}",
                                },
                            },
                        ],
                    }
                ],
                model="llama-3.2-11b-vision-preview",
            )
            return chat_completion.choices[0].message.content
        except Exception as e:
            print(f"Vision Error: {e}")
            return "[Error analyzing image]"
